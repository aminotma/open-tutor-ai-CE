"""
Context Retrieval Engine Router

This module implements the Context Retrieval Engine that combines:
- Pedagogical documents (RAG)
- Internal memory (episodic, semantic, procedural, behavioral)
- Generated summaries

Results are filtered and ranked according to pedagogical relevance.
"""

from enum import Enum
from typing import List, Optional, Dict, Any, cast
from uuid import uuid4
from dataclasses import dataclass
from datetime import datetime, timezone, timedelta
from math import exp
import difflib
import json
import os
import re
import tiktoken
from pathlib import Path

import chromadb
from chromadb.api.shared_system_client import SharedSystemClient
from chromadb.config import Settings

from fastapi import APIRouter, Depends, HTTPException, Query
from pydantic import BaseModel, Field

from open_webui.internal.db import get_db, engine
from open_webui.utils.auth import get_verified_user
from open_tutorai.config import CONTEXT_RETRIEVAL_CONFIG, get_openai_api_key, get_openai_base_url
from open_tutorai.models.database import Memory
from sqlalchemy.orm import sessionmaker

from langchain_chroma import Chroma
from langchain_core.prompts import ChatPromptTemplate
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_core.documents import Document as LangChainDocument
from langchain_openai import ChatOpenAI

router = APIRouter(tags=["context"])


def get_db_session():
    """Get a database session using the same engine as OpenWebUI"""
    Session = sessionmaker(bind=engine)
    return Session()


# ============================================================================
# CHROMADB SETUP
# ============================================================================

# Initialize ChromaDB client
# Imports LangChain

# Client ChromaDB — initialisé en différé pour éviter les crashs à l'import
_chroma_client = None

# Collection for pedagogical documents
DOCUMENTS_COLLECTION = "pedagogical_documents"


def _get_chroma_client():
    """Retourne (et crée si besoin) le client ChromaDB partagé."""
    global _chroma_client
    if _chroma_client is None:
        # Chemin absolu : stable peu importe le répertoire de travail du processus
        db_path = Path(__file__).resolve().parents[2] / "data" / "vector_db"
        db_path.mkdir(parents=True, exist_ok=True)
        settings = Settings(anonymized_telemetry=False)
        try:
            _chroma_client = chromadb.PersistentClient(
                path=str(db_path),
                settings=settings,
            )
        except ValueError as e:
            error_message = str(e)
            if "An instance of Chroma already exists for" in error_message:
                SharedSystemClient.clear_system_cache()
                _chroma_client = chromadb.PersistentClient(
                    path=str(db_path),
                    settings=settings,
                )
            else:
                raise
    return _chroma_client


def get_or_create_collection(collection_name: str):
    """Récupère ou crée une collection ChromaDB."""
    client = _get_chroma_client()
    try:
        return client.get_collection(name=collection_name)
    except Exception:
        return client.create_collection(name=collection_name)

# Fonction d'embeddings partagée entre indexation ET recherche
# Utilise le même modèle qu'avant (all-MiniLM-L6-v2) mais géré par LangChain
_embedding_function = None

def get_embedding_function() -> HuggingFaceEmbeddings:
    """Singleton pour éviter de recharger le modèle à chaque appel."""
    global _embedding_function
    if _embedding_function is None:
        _embedding_function = HuggingFaceEmbeddings(
            model_name="all-MiniLM-L6-v2",
            model_kwargs={"device": "cpu"},
            encode_kwargs={"normalize_embeddings": True}
        )
    return _embedding_function

# Vectorstore LangChain — utilise le même client ChromaDB
_vectorstore = None

def get_vectorstore() -> Chroma:
    """Retourne le vectorstore LangChain."""
    global _vectorstore
    if _vectorstore is None:
        _vectorstore = Chroma(
            client=_get_chroma_client(),
            collection_name=DOCUMENTS_COLLECTION,
            embedding_function=get_embedding_function(),
        )
    return _vectorstore


# ============================================================================
# DOCUMENT INDEXING FUNCTIONS
# ============================================================================

def index_document_to_chromadb(
    doc_id: str,
    content: str,
    metadata: Dict[str, Any],
    collection_name: str = DOCUMENTS_COLLECTION
) -> bool:
    """
    Indexe via LangChain pour garantir la cohérence des embeddings
    entre l'indexation et la recherche.
    """
    try:
        from langchain_core.documents import Document as LangChainDocument

        # Nettoyage des métadonnées (LangChain accepte str/int/float/bool)
        clean_metadata = {}
        for key, value in metadata.items():
            if isinstance(value, (str, int, float, bool)):
                clean_metadata[key] = value
            else:
                clean_metadata[key] = json.dumps(value)

        # Store the document id in metadata so we can recover it during retrieval
        clean_metadata["doc_id"] = doc_id

        doc = LangChainDocument(
            page_content=content,
            metadata=clean_metadata
        )

        vectorstore = get_vectorstore()

        # Vérifie si déjà indexé via le client bas niveau
        try:
            existing = _get_chroma_client().get_collection(
                DOCUMENTS_COLLECTION
            ).get(ids=[doc_id])
            if existing and existing.get("ids"):
                return True  # Déjà indexé, on ne réindexe pas
        except Exception:
            pass  # Collection vide ou doc absent, on continue

        vectorstore.add_documents(
            documents=[doc],
            ids=[doc_id]
        )
        return True

    except Exception as e:
        print(f"Error indexing document {doc_id}: {str(e)}")
        return False


def _extract_text_from_file(file_path: str) -> Optional[str]:
    """Extract text from common document formats for indexing."""
    path = Path(file_path)
    suffix = path.suffix.lower()

    try:
        if suffix == ".pdf":
            try:
                from pypdf import PdfReader
            except ImportError:
                from PyPDF2 import PdfReader

            reader = PdfReader(str(path))
            text_parts = []
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
            return "\n\n".join(text_parts).strip()

        if suffix == ".docx":
            try:
                import docx2txt
                return docx2txt.process(str(path)) or ""
            except Exception:
                pass

        if suffix == ".pptx":
            try:
                from pptx import Presentation
                from pptx.shapes.base import BaseShape
                
                presentation = Presentation(str(path))
                slides = []
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        shape_any = cast(Any, shape)
                        if getattr(shape_any, "has_text_frame", False):
                            text = getattr(shape_any, "text", None)
                            if text:
                                slides.append(text)
                return "\n\n".join(slides).strip()
            except Exception:
                pass

        if suffix in {".txt", ".md", ".json", ".csv", ".py", ".yaml", ".yml"}:
            return path.read_text(encoding="utf-8", errors="ignore")

        # Fallback: try reading as UTF-8 text
        return path.read_text(encoding="utf-8", errors="ignore")

    except Exception:
        try:
            with open(path, "rb") as f:
                return f.read().decode("utf-8", errors="ignore")
        except Exception:
            return None


def index_uploaded_document_to_chromadb(file_path: str, user_id: str, title: str) -> bool:
    """
    Indexe un document uploadé par l'utilisateur dans ChromaDB.
    
    Args:
        file_path: Chemin absolu vers le fichier à indexer
        user_id: ID de l'utilisateur qui a uploadé le fichier
        title: Titre du document pour les métadonnées
    
    Returns:
        True si l'indexation a réussi, False sinon
    """
    try:
        from langchain_core.documents import Document as LangChainDocument
        from pathlib import Path

        path = Path(file_path)
        if not path.exists():
            print(f"File not found: {file_path}")
            return False

        content = _extract_text_from_file(str(path))
        if not content:
            print(f"No text extracted from file: {file_path}")
            return False

        doc_id = f"upload_{uuid4()}"

        metadata = {
            "type": "uploaded_document",
            "source": "user_upload",
            "user_id": user_id,
            "title": title,
            "file_name": path.name,
            "file_path": str(path),
            "file_type": path.suffix.lower().lstrip('.'),
            "indexed_at": datetime.now(timezone.utc).timestamp(),
            "doc_id": doc_id
        }

        clean_metadata = {}
        for key, value in metadata.items():
            if isinstance(value, (str, int, float, bool)):
                clean_metadata[key] = value
            else:
                clean_metadata[key] = json.dumps(value)

        doc = LangChainDocument(
            page_content=content,
            metadata=clean_metadata
        )

        vectorstore = get_vectorstore()
        vectorstore.add_documents(
            documents=[doc],
            ids=[doc_id]
        )

        print(f"Successfully indexed document: {title} (ID: {doc_id})")
        return True

    except Exception as e:
        print(f"Error indexing uploaded document {file_path}: {str(e)}")
        return False


def index_local_documents_to_chromadb(collection_name: str = DOCUMENTS_COLLECTION) -> int:
    """Index pedagogical documents from the docs/ directory into ChromaDB.
    Only the docs/ folder is searched — backend source files are never indexed."""
    repo_root = Path(__file__).resolve().parents[3]
    search_paths = [repo_root / "docs"]  # backend/ excluded: contains only developer docs
    # Répertoires contenant des données utilisateur — jamais indexés
    EXCLUDED_DIRS = {"data/cache", "data/uploads", "data/vector_db", "__pycache__", ".git", "node_modules"}
    # Fichiers techniques non pédagogiques — jamais indexés
    EXCLUDED_NAMES = {"requirements.txt", "requirements-dev.txt", "package.json", "package-lock.json",
                      ".env", ".env.example", "CHANGELOG.md", "CONTRIBUTING.md", "SECURITY.md"}
    supported_exts = {".md", ".txt", ".pdf"}  # JSON exclus : fichiers de session/config
    indexed_count = 0
    collection = get_or_create_collection(collection_name)

    for base_path in search_paths:
        if not base_path.exists():
            continue

        for file_path in base_path.rglob("*"):
            if not file_path.is_file() or file_path.suffix.lower() not in supported_exts:
                continue

            # Exclure tout chemin contenant un répertoire interdit ou un nom de fichier technique
            relative = file_path.relative_to(repo_root).as_posix()
            if any(excl in relative for excl in EXCLUDED_DIRS):
                continue
            if file_path.name in EXCLUDED_NAMES:
                continue

            doc_id = f"local_{file_path.relative_to(repo_root).as_posix()}"
            try:
                existing = collection.get(ids=[doc_id])
                if existing and existing.get("ids"):
                    continue
            except Exception:
                pass

            content = _extract_text_from_file(str(file_path))
            if not content:
                continue

            metadata = {
                "type": "local_document",
                "source": "repository",
                "title": file_path.name,
                "path": str(file_path.relative_to(repo_root)),
                "file_type": file_path.suffix.lower().lstrip('.'),
                "indexed_at": datetime.now(timezone.utc).timestamp(),
                "doc_id": doc_id
            }

            if index_document_to_chromadb(doc_id, content, metadata, collection_name=collection_name):
                indexed_count += 1

    return indexed_count


# ============================================================================
# PYDANTIC MODELS
# ============================================================================

class ContextRetrievalRequest(BaseModel):
    """Request model for context retrieval"""
    query: str = Field(..., min_length=1, description="Search query")
    max_results: int = Field(20, ge=1, le=50, description="Maximum results to return")
    include_source_types: Optional[List[str]] = Field(
        None,
        description="Filter by source types: 'pedagogical', 'memory', 'summary'"
    )
    memory_types: Optional[List[str]] = Field(
        None,
        description="Filter memories by types: 'episodic', 'semantic', 'procedural', 'behavioral'"
    )
    pedagogical_level: Optional[str] = Field(
        "intermediate",
        description="User pedagogical level: 'beginner', 'intermediate', 'advanced'"
    )
    learning_objectives: Optional[List[str]] = Field(
        None,
        description="Learning objectives to align retrieved context with user goals"
    )
    topic: Optional[str] = Field(
        None,
        description="Current course/topic — isolates memories to this course. "
                    "Empty string or None retrieves the 'general' bucket (no course active)."
    )


class DocumentInteractionRequest(BaseModel):
    """Request model to record user interactions with a pedagogical document."""
    interaction_type: Optional[str] = Field(
        None,
        description="Optional type of user interaction, e.g. 'view', 'select', 'answer_question'"
    )
    interaction_increment: int = Field(
        1,
        ge=1,
        description="Number of interaction events to add to the document engagement count"
    )
    answered_questions_increment: Optional[int] = Field(
        0,
        ge=0,
        description="Optional number of answered questions to increment for this document"
    )
    last_updated: Optional[float] = Field(
        None,
        description="Optional timestamp to set as last interaction time"
    )


class ScoresSchema(BaseModel):
    """Score breakdown for context item"""
    relevance: float = Field(..., description="Relevance score (0-1)")
    engagement: float = Field(..., description="Engagement score (0-1)")
    recency: float = Field(..., description="Recency score (0-1)")
    user_alignment: float = Field(..., description="User preference alignment (0-1)")
    composite: float = Field(..., description="Composite score (0-1)")
    normalized: float = Field(..., description="Normalized score within result set (0-1)")


class MetadataSchema(BaseModel):
    """Metadata for context item"""
    type: str = Field(..., description="Content type")
    created_at: Optional[float] = Field(None, description="Creation timestamp")
    last_updated: Optional[float] = Field(None, description="Last update timestamp")
    source_id: str = Field(..., description="Source identifier")
    title: Optional[str] = Field(None, description="Title if applicable")
    subject_domain: Optional[str] = Field(None, description="Subject domain")
    source: Optional[str] = Field(None, description="Source of the pedagogical document")
    interaction_count: Optional[int] = Field(None, description="Number of user interactions recorded for this document")
    answered_questions: Optional[int] = Field(None, description="Number of answered questions associated with this document")


class ContextRetrievalResponse(BaseModel):
    """Response model for a single context item"""
    rank: int = Field(..., description="Ranking position")
    source: str = Field(..., description="Source type: pedagogical, memory, summary")
    source_id: str = Field(..., description="Source identifier")
    content_preview: str = Field(..., description="Preview of content (first 300 chars)")
    full_content: str = Field(..., description="Full content")
    metadata: MetadataSchema = Field(..., description="Content metadata")
    scores: ScoresSchema = Field(..., description="Score breakdown")


# ============================================================================
# DATACLASSES FOR INTERNAL PROCESSING
# ============================================================================

@dataclass
class NormalizedContextItem:
    """Normalized context item from any source"""
    source_type: str  # 'pedagogical', 'memory', 'summary'
    id: str
    content: str
    metadata: Dict[str, Any]
    raw_score: float


# ============================================================================
# SUMMARIZATION LAYER
# ============================================================================

def _summarize_mechanical(content: str, max_tokens: int = 500) -> str:
    """Fallback: token-aware extraction of first/middle/last sentences."""
    if not content:
        return ""
    try:
        enc = tiktoken.get_encoding("cl100k_base")
        tokens = enc.encode(content)
        if len(tokens) <= max_tokens:
            return content
    except Exception:
        if len(content) <= max_tokens * 4:
            return content
        return content[:max_tokens * 4] + "..."

    sentences = [s.strip() for s in content.split('.') if s.strip()]
    if len(sentences) <= 3:
        return enc.decode(tokens[:max_tokens]) + "..."

    key_sentences = [sentences[0]]
    middle_idx = len(sentences) // 2
    if middle_idx not in (0, len(sentences) - 1):
        key_sentences.append(sentences[middle_idx])
    if len(sentences) > 1:
        key_sentences.append(sentences[-1])

    summary = '. '.join(key_sentences)
    summary_tokens = enc.encode(summary)
    if len(summary_tokens) <= max_tokens:
        return summary
    return enc.decode(summary_tokens[:max_tokens - 10]) + "..."


_SUMMARIZE_PROMPT = ChatPromptTemplate.from_messages([
    (
        "system",
        "You are a pedagogical context summarizer. "
        "Summarize the content below as concisely as possible, focusing on what is most "
        "relevant to the learner's query. Return only the summary text — no preamble, no headings."
    ),
    (
        "human",
        "QUERY: {query}\n\nCONTENT:\n{content}"
    ),
])


async def summarize_interactions(content: str, query: str = "", max_tokens: int = 500) -> str:
    """
    LLM-based summarizer: generates a coherent, query-relevant summary.
    Falls back to mechanical sentence extraction on any error.
    """
    if not content:
        return ""

    # Skip LLM call when content already fits
    try:
        enc = tiktoken.get_encoding("cl100k_base")
        if len(enc.encode(content)) <= max_tokens:
            return content
    except Exception:
        if len(content) <= max_tokens * 4:
            return content

    lc_cfg = CONTEXT_RETRIEVAL_CONFIG.get("langchain", {})
    llm = ChatOpenAI(
        model=lc_cfg.get("llm_model", "gpt-4o-mini"),
        temperature=0.0,
        api_key=get_openai_api_key(),
        base_url=get_openai_base_url(),
    )
    chain = _SUMMARIZE_PROMPT | llm

    try:
        result = await chain.ainvoke({
            "query": query or "general context",
            "content": content[:2000],
        })
        summary = result.content.strip()
        # Enforce token budget on the LLM output
        try:
            enc = tiktoken.get_encoding("cl100k_base")
            toks = enc.encode(summary)
            if len(toks) > max_tokens:
                return enc.decode(toks[:max_tokens]) + "..."
        except Exception:
            pass
        return summary
    except Exception:
        return _summarize_mechanical(content, max_tokens)


def sliding_window_filter(
    items: List[Dict],
    window_size: int = 10,
    score_threshold: float = 0.3
) -> List[Dict]:
    """
    Apply sliding window to keep only the most relevant recent items
    
    Keeps top N items by score, preferring more recent ones
    """
    if len(items) <= window_size:
        return items
    
    # Sort by composite score (descending) then by recency
    sorted_items = sorted(
        items,
        key=lambda x: (x.get("composite_score", 0), x.get("recency_score", 0)),
        reverse=True
    )
    
    # Keep only top window_size items above threshold
    filtered = []
    for item in sorted_items[:window_size]:
        if item.get("composite_score", 0) >= score_threshold:
            filtered.append(item)
    
    return filtered


def extract_key_elements(content: str, query: str) -> str:
    """
    Extract key elements from content based on query relevance
    
    Identifies sentences most relevant to the query
    """
    if not content or not query:
        return content

    sentences = [s.strip() for s in content.split('.') if s.strip()]
    if not sentences:
        return content

    query_terms = set(query.lower().split())
    
    # Score each sentence by query term matches
    scored_sentences = []
    for sentence in sentences:
        sentence_lower = sentence.lower()
        matches = sum(1 for term in query_terms if term in sentence_lower)
        score = matches / len(query_terms) if query_terms else 0
        scored_sentences.append((sentence, score))
    
    # Keep sentences with score > 0, sorted by score
    relevant_sentences = [
        sentence for sentence, score in scored_sentences
        if score > 0
    ]
    
    if not relevant_sentences:
        # Fallback: keep first and last sentences
        relevant_sentences = sentences[:1] + sentences[-1:] if len(sentences) > 1 else sentences
    
    # Limit to 3 key sentences max
    relevant_sentences = relevant_sentences[:3]
    
    return '. '.join(relevant_sentences)


def forget_irrelevant_sentences(
    content: str,
    query: str,
    min_match_ratio: float = 0.15
) -> str:
    """
    Remove sentences from content that do not match the query well.
    """
    if not content or not query:
        return content

    sentences = [s.strip() for s in content.split('.') if s.strip()]
    if not sentences:
        return content

    query_terms = set(query.lower().split())
    if not query_terms:
        return content

    kept_sentences = []
    for sentence in sentences:
        sentence_lower = sentence.lower()
        matches = sum(1 for term in query_terms if term in sentence_lower)
        match_ratio = matches / len(query_terms)
        if match_ratio >= min_match_ratio:
            kept_sentences.append(sentence)

    if not kept_sentences:
        return content

    return '. '.join(kept_sentences)


def forget_irrelevant_context_items(
    items: List[Dict],
    threshold: float = 0.15,
    min_relevance: float = 0.1
) -> List[Dict]:
    """
    Remove context items that are unlikely to contribute to a useful summary.
    """
    if not items:
        return items

    kept = []
    for item in items:
        if item.get("normalized_score", 0) >= threshold:
            kept.append(item)
            continue

        if item.get("relevance_score", 0) >= min_relevance:
            kept.append(item)
            continue

    return kept if kept else items[:1]


async def apply_summarization_layer(
    ranked_items: List[Dict],
    query: str,
    config: Dict
) -> List[Dict]:
    """
    Apply summarization layer to reduce context size while preserving essential information
    
    Techniques:
    - Forget irrelevant context items
    - Remove irrelevant sentences from long content
    - Summarize long interactions
    - Apply sliding window filtering
    - Extract key elements based on query
    """
    summarization_config = config.get("summarization", {})
    max_content_length = summarization_config.get("max_content_length", 1000)
    window_size = summarization_config.get("sliding_window_size", 10)
    score_threshold = summarization_config.get("score_threshold", 0.3)

    forget_config = summarization_config.get("forget_irrelevant", {})
    forget_enabled = forget_config.get("enabled", True)
    forget_threshold = forget_config.get("context_threshold", 0.15)
    forget_min_relevance = forget_config.get("min_relevance", 0.1)
    sentence_match_ratio = forget_config.get("sentence_match_ratio", 0.15)
    
    summarized_items = []
    
    for item in ranked_items:
        summarized_item = item.copy()
        content = item.get("content", "")
        
        # Remove irrelevant sentences first to reduce noise
        if forget_enabled and len(content) > max_content_length:
            content = forget_irrelevant_sentences(
                content,
                query,
                min_match_ratio=sentence_match_ratio
            )

        # Apply summarization techniques
        if len(content) > max_content_length:
            # Extract key elements first
            content = extract_key_elements(content, query)
            
            # Then summarize if still too long
            if len(content) > max_content_length:
                content = await summarize_interactions(content, query=query, max_tokens=max_content_length // 4)
        
        summarized_item["content"] = content
        summarized_item["original_length"] = len(item.get("content", ""))
        summarized_item["summarized_length"] = len(content)
        
        summarized_items.append(summarized_item)
    
    if forget_enabled:
        summarized_items = forget_irrelevant_context_items(
            summarized_items,
            threshold=forget_threshold,
            min_relevance=forget_min_relevance
        )

    # Apply sliding window to final results
    final_items = sliding_window_filter(
        summarized_items,
        window_size=window_size,
        score_threshold=score_threshold
    )
    
    return final_items


# ============================================================================
# STEP 1: RETRIEVE FROM MULTIPLE SOURCES
# ============================================================================

def _read_text_file(path: Path) -> Optional[str]:
    try:
        text = path.read_text(encoding="utf-8", errors="ignore")
        return text
    except Exception:
        return None


def _tokenize_text(text: str) -> List[str]:
    return [token for token in re.findall(r"\w{3,}", text.lower())]


def _calculate_text_similarity(query: str, content: str) -> float:
    query_tokens = set(_tokenize_text(query))
    content_tokens = set(_tokenize_text(content))
    if not query_tokens or not content_tokens:
        return 0.0

    overlap = len(query_tokens & content_tokens) / max(1, len(query_tokens | content_tokens))
    sequence_ratio = difflib.SequenceMatcher(None, query.lower(), content.lower()).ratio()
    return min(1.0, 0.6 * overlap + 0.4 * sequence_ratio)


def _score_document_relevance(content: str, query: str) -> float:
    query_terms = re.findall(r"\w+", query.lower())
    if not query_terms:
        return 0.0

    content_lower = content.lower()
    matches = sum(1 for term in query_terms if term in content_lower)
    exact_score = min(1.0, matches / max(len(query_terms), 1))
    similarity_score = _calculate_text_similarity(query, content)
    return float(max(exact_score, similarity_score))

async def retrieve_pedagogical_documents(
    user_id: str,
    query: str,
    top_k: int = 5
) -> List[Dict]:
    """
    Récupère les documents via LangChain similarity_search_with_score.
    Retourne le même format de dictionnaire que l'ancienne version
    pour ne pas casser les consommateurs existants.
    """
    try:
        if not query or not query.strip():
            return []

        # Indexation lazy des documents locaux (inchangée)
        indexed_count = index_local_documents_to_chromadb()
        if indexed_count > 0:
            print(f"Indexed {indexed_count} new local documents")

        vectorstore = get_vectorstore()

        # Filtre : documents appartenant à cet utilisateur
        # OU documents sans propriétaire (ressources pédagogiques partagées)
        doc_filter = {
            "$or": [
                {"user_id": {"$eq": user_id}},
                {"user_id": {"$eq": ""}},
            ]
        }

        # Recherche sémantique via LangChain
        # Retourne List[Tuple[LangChainDocument, float]]
        results_with_scores = vectorstore.similarity_search_with_score(
            query=query,
            k=top_k,
            filter=doc_filter,
        )

        documents = []
        for lc_doc, distance in results_with_scores:
            # ChromaDB retourne une distance cosinus (0=identique, 2=opposé)
            # On convertit en score de similarité [0, 1]
            vector_score = max(0.0, 1.0 - (distance / 2.0))

            metadata = lc_doc.metadata or {}

            # Extraction de l'id (stocké dans les métadonnées ou généré)
            doc_id = metadata.get("doc_id", metadata.get("path", f"doc_{hash(lc_doc.page_content)}"))

            documents.append({
                "id": doc_id,
                "source_type": "rag",
                "content": lc_doc.page_content,
                "metadata": {
                    "type": metadata.get("type", "document"),
                    "title": metadata.get("title", doc_id),
                    "path": metadata.get("path", ""),
                    "source": metadata.get("source", "unknown"),
                    "file_type": metadata.get("file_type", ""),
                    "user_id": metadata.get("user_id", ""),
                    "indexed_at": metadata.get("indexed_at", ""),
                    "uploaded_at": metadata.get("uploaded_at", ""),
                    "last_updated": metadata.get("last_updated"),
                    "interaction_count": metadata.get("interaction_count", 0),
                    "answered_questions": metadata.get("answered_questions", 0),
                    "doc_id": doc_id
                },
                "relevance_score": vector_score,
                "vector_score": vector_score
            })

        return documents

    except Exception as e:
        print(f"Error retrieving pedagogical documents (LangChain): {str(e)}")
        return await _retrieve_pedagogical_documents_fallback(user_id, query, top_k)

async def _retrieve_pedagogical_documents_fallback(
    user_id: str,
    query: str,
    top_k: int = 5
) -> List[Dict]:
    """
    Fallback retrieval using file-walk when ChromaDB is unavailable
    
    This maintains backward compatibility and ensures the system works
    even if vector search fails.
    """
    try:
        if not query or not query.strip():
            return []

        repo_root = Path(__file__).resolve().parents[4]
        search_paths = [repo_root / "docs", repo_root / "backend"]
        documents = []

        for base_path in search_paths:
            if not base_path.exists():
                continue
            for root, _, files in os.walk(base_path):
                for file_name in files:
                    if not file_name.lower().endswith((".md", ".txt", ".json")):
                        continue
                    file_path = Path(root) / file_name
                    text = _read_text_file(file_path)
                    if not text:
                        continue

                    relevance_score = _score_document_relevance(text, query)
                    if relevance_score <= 0:
                        continue

                    documents.append({
                        "id": str(file_path.relative_to(repo_root)),
                        "source_type": "rag",
                        "content": text,
                        "metadata": {
                            "title": file_name,
                            "path": str(file_path.relative_to(repo_root)),
                            "source": "local_document",
                        },
                        "relevance_score": relevance_score,
                    })

        documents.sort(key=lambda x: x["relevance_score"], reverse=True)
        return documents[:top_k]

    except Exception as e:
        print(f"Error in fallback retrieval: {str(e)}")
        return []

async def retrieve_pedagogical_documents_as_langchain(
    user_id: str,
    query: str,
    top_k: int = 5
) -> List["LangChainDocument"]:
    """
    Version qui retourne des objets LangChain Document natifs.
    Utilisée directement par les RetrievalChain LangChain.
    """
    from langchain_core.documents import Document as LangChainDocument

    try:
        if not query or not query.strip():
            return []

        vectorstore = get_vectorstore()
        return vectorstore.similarity_search(query=query, k=top_k)

    except Exception as e:
        print(f"Error in LangChain retrieval: {str(e)}")
        return []

async def retrieve_internal_memory(
    user_id: str,
    query: str,
    memory_types: Optional[List[str]] = None,
    limit: int = 10,
    db = None,
    topic: Optional[str] = None,
) -> List[Dict]:
    """
    Retrieve internal memories matching the query with intelligent summarization and selection

    Source: opentutorai_memory table
    Uses: semantic relevance scoring with memory type prioritization and content summarization.
    """
    if db is None:
        return []

    try:
        db_query = db.query(Memory).filter(Memory.user_id == user_id)
        if memory_types and len(memory_types) > 0:
            db_query = db_query.filter(Memory.memory_type.in_(memory_types))

        # Retrieve broader set of candidate memories
        memories = []
        if query and query.strip():
            search_pattern = f"%{query}%"
            candidate_query = db_query.filter(Memory.content.ilike(search_pattern))
            memories = candidate_query.order_by(
                Memory.updated_at.desc().nullslast(),
                Memory.created_at.desc()
            ).limit(limit * 5).all()

        if not memories:
            memories = db_query.order_by(
                Memory.updated_at.desc().nullslast(),
                Memory.created_at.desc()
            ).limit(limit * 15).all()

        # Topic isolation: keep only memories that belong to the current course bucket.
        # topic=None or "" → "general" bucket (memories stored without a course).
        # topic="Python"   → only memories tagged with topic="Python".
        target_topic = (topic or "").strip().lower()
        memories = [
            m for m in memories
            if (m.memory_metadata or {}).get("topic", "").strip().lower() == target_topic
        ]

        scored_memories = []
        for memory in memories:
            semantic_score = (
                calculate_relevance(memory.content, query)
                if query and query.strip()
                else 0.5
            )

            # Skip very low relevance memories
            if query and query.strip() and semantic_score <= 0.05:
                continue

            # Calculate importance score based on memory type and recency
            type_importance = {
                "episodic": 1.0,    # Learning experiences
                "semantic": 1.2,    # Learned concepts
                "procedural": 1.1,  # Methods and strategies
                "behavioral": 0.9   # Session summaries
            }

            base_importance = type_importance.get(memory.memory_type, 0.8)

            # Recency boost (newer memories are more important)
            recency_score = 0.5
            if memory.updated_at:
                days_since_update = (datetime.now(timezone.utc) - memory.updated_at).days
                recency_score = max(0.1, 1.0 - (days_since_update / 30.0))  # Decay over 30 days

            # Content quality score (longer, more detailed memories are better)
            content_length = len(memory.content) if memory.content else 0
            quality_score = min(1.0, content_length / 500.0)  # Max at 500 chars

            # Final importance score
            importance_score = (base_importance * 0.4) + (recency_score * 0.3) + (quality_score * 0.3)

            # Combined relevance score
            combined_score = (semantic_score * 0.7) + (importance_score * 0.3)

            scored_memories.append({
                "memory": memory,
                "semantic_score": semantic_score,
                "importance_score": importance_score,
                "combined_score": combined_score
            })

        # Sort by combined score (descending)
        scored_memories.sort(key=lambda x: x["combined_score"], reverse=True)

        results = []
        for item in scored_memories[:limit]:
            memory = item["memory"]

            # Summarize long content to keep context manageable
            content = memory.content
            if len(content) > 300:
                content = await summarize_interactions(content, query=query, max_tokens=100)

            results.append({
                "id": memory.id,
                "type": memory.memory_type,
                "content": content,
                "metadata": memory.memory_metadata or {},
                "created_at": memory.created_at.timestamp() if memory.created_at else None,
                "updated_at": memory.updated_at.timestamp() if memory.updated_at else None,
                "textual_score": item["semantic_score"],
                "importance_score": item["importance_score"],
                "combined_score": item["combined_score"]
            })

        return results

    except Exception as e:
        print(f"Error retrieving internal memory: {str(e)}")
        return []


async def retrieve_generated_summaries(
    user_id: str,
    query: str,
    cache_ttl_hours: int = 24,
    limit: int = 5,
    topic: Optional[str] = None,
) -> List[Dict]:
    """
    Retrieve generated summaries from cache
    
    Source: backend/data/cache/summaries/
    Uses cache with TTL
    """
    try:
        cache_dir = Path("backend/data/cache/summaries")
        cache_dir.mkdir(parents=True, exist_ok=True)
        
        candidates = []
        current_time = datetime.now(timezone.utc)

        # Scan cache directory — files are named {user_id}_{chat_id}_{ts}.json
        # Read a wider pool (limit × 4) so relevance-based re-ranking has enough
        # candidates, without reading every historical file for heavy users.
        if cache_dir.exists():
            user_files = sorted(
                cache_dir.glob(f"{user_id}_*.json"),
                key=lambda p: p.stat().st_mtime,
                reverse=True,
            )[:limit * 4]

            for cache_file in user_files:
                try:
                    with open(cache_file, 'r', encoding='utf-8') as f:
                        cached_summary = json.load(f)

                    # Check TTL
                    created_timestamp = cached_summary.get("created_at", 0)
                    created_date = datetime.fromtimestamp(created_timestamp, tz=timezone.utc)
                    age = (current_time - created_date).total_seconds() / 3600

                    if age <= cache_ttl_hours:
                        # Topic isolation — same bucket logic as retrieve_internal_memory
                        file_topic = cached_summary.get("topic", "").strip().lower()
                        target_topic = (topic or "").strip().lower()
                        if file_topic != target_topic:
                            continue

                        relevance = calculate_relevance(cached_summary.get("text", ""), query)
                        if relevance > 0.1:
                            candidates.append({
                                "id": cached_summary.get("id", cache_file.stem),
                                "text": cached_summary.get("text", ""),
                                "source_conversation": cached_summary.get("source_conversation"),
                                "created_at": created_timestamp,
                                "summary_score": relevance,
                            })

                except (json.JSONDecodeError, IOError):
                    continue

        # Return the most relevant summaries, not just the most recent ones
        candidates.sort(key=lambda x: x["summary_score"], reverse=True)
        return candidates[:limit]
        
    except Exception as e:
        print(f"Error retrieving generated summaries: {str(e)}")
        return []


# ============================================================================
# STEP 2: NORMALIZE AND STANDARDIZE
# ============================================================================

def normalize_context(
    documents: List[Dict],
    memories: List[Dict],
    summaries: List[Dict]
) -> List[NormalizedContextItem]:
    """
    Normalize results from all sources to unified format
    """
    normalized = []
    
    # Normalize pedagogical documents
    for doc in documents:
        content = doc.get("content", "")
        if len(content) > 5000:
            content = content[:5000] + "..."
        
        doc_meta = doc.get("metadata", {})
        normalized.append(NormalizedContextItem(
            source_type="pedagogical",
            id=doc.get("id", ""),
            content=content,
            metadata={
                "type": doc_meta.get("type", "document"),
                "title": doc_meta.get("title", doc.get("id", "")),
                "created_at": doc_meta.get("created_at"),
                "last_updated": doc_meta.get("last_updated"),
                "source_id": doc.get("id", ""),
                "source": doc_meta.get("source", ""),
                "interaction_count": doc_meta.get("interaction_count", 0),
                "answered_questions": doc_meta.get("answered_questions", 0)
            },
            raw_score=doc.get("vector_score", 0.0)
        ))

    # Normalize memories
    for memory in memories:
        content = memory.get("content", "")
        if len(content) > 5000:
            content = content[:5000] + "..."
        
        normalized.append(NormalizedContextItem(
            source_type="memory",
            id=memory.get("id", ""),
            content=content,
            metadata={
                "type": memory.get("type", "semantic"),
                "created_at": memory.get("created_at"),
                "last_updated": memory.get("updated_at"),
                "source_id": memory.get("id", "")
            },
            raw_score=memory.get("textual_score", 0.0)
        ))
    
    # Normalize summaries
    for summary in summaries:
        content = summary.get("text", "")
        if len(content) > 5000:
            content = content[:5000] + "..."
        
        normalized.append(NormalizedContextItem(
            source_type="summary",
            id=summary.get("id", ""),
            content=content,
            metadata={
                "type": "summary",
                "source_conversation": summary.get("source_conversation"),
                "created_at": summary.get("created_at"),
                "source_id": summary.get("id", "")
            },
            raw_score=summary.get("summary_score", 0.0)
        ))
    
    return normalized


# ============================================================================
# STEP 3: ENRICH WITH PEDAGOGICAL METADATA
# ============================================================================

def calculate_relevance(content: str, query: str) -> float:
    """
    Calculate textual relevance score
    
    Returns: score between 0 and 1
    """
    if not query or not content:
        return 0.0
    
    query_lower = query.lower()
    content_lower = content.lower()
    
    # Count query terms found
    terms = [t.strip() for t in query_lower.split() if t.strip()]
    if not terms:
        return 0.0
    
    matches = sum(1 for term in terms if term in content_lower)
    exact_score = min(matches / len(terms), 1.0)
    similarity_score = _calculate_text_similarity(query, content)
    return max(exact_score, similarity_score)


def deduce_pedagogical_level(content: str) -> str:
    """
    Deduce pedagogical level from content characteristics
    
    Returns: 'beginner', 'intermediate', or 'advanced'
    """
    if not content:
        return "intermediate"
    
    content_length = len(content)
    
    # Heuristic based on content length and complexity
    if content_length < 100:
        return "beginner"
    elif content_length < 500:
        return "intermediate"
    else:
        return "advanced"


def calculate_recency_score(created_at: Optional[float]) -> float:
    """
    Calculate recency score using exponential decay
    
    Returns: score between 0 and 1, 1 = very recent, 0 = very old
    """
    if created_at is None:
        return 0.5
    
    current_time = datetime.now(timezone.utc).timestamp()
    age_seconds = current_time - created_at
    age_days = age_seconds / (24 * 3600)
    
    # Exponential decay with 30-day half-life
    recency = exp(-age_days / 30)
    
    return min(max(recency, 0.0), 1.0)


def calculate_engagement_score(item_type: str, metadata: Dict) -> float:
    """
    Calculate engagement score based on item type and metadata
    
    Returns: score between 0 and 1
    """
    score = 0.0
    
    # Episodic memories get higher engagement
    if item_type == "memory" and metadata.get("type") == "episodic":
        score += 0.3

    # Summaries — baseline engagement (content was relevant enough to be summarised)
    if item_type == "summary":
        score += 0.2

    # Pedagogical documents should have a baseline engagement
    # because learners are interacting with them or answering questions.
    if item_type == "pedagogical":
        score += 0.2

        # Boost if we know the content was uploaded or indexed by the learner
        if metadata.get("type") in {"uploaded_document", "local_document", "document"}:
            score += 0.1

    # Recently updated items get bonus
    if metadata.get("last_updated"):
        score += 0.2

    # Support explicit interaction metadata if available
    try:
        answered = float(metadata.get("answered_questions", 0))
        score += min(answered * 0.05, 0.3)
    except (TypeError, ValueError):
        pass

    try:
        interactions = float(metadata.get("interaction_count", 0))
        score += min(interactions * 0.03, 0.3)
    except (TypeError, ValueError):
        pass

    return min(score, 1.0)


def calculate_user_alignment(
    subject_domain: str,
    user_profile: Dict,
    content: str
) -> float:
    """
    Calculate user preference alignment score based on subject domain and learning objectives.
    """
    user_interests = [str(item).lower() for item in user_profile.get("interests", []) if item]
    learning_objectives = [str(item).lower() for item in user_profile.get("learning_objectives", []) if item]
    subject = str(subject_domain or "").lower()
    text = str(content or "").lower()

    if subject and subject in user_interests:
        return 1.0

    if any(obj in subject for obj in learning_objectives):
        return 0.9

    if any(obj in text for obj in learning_objectives):
        return 0.8

    return 0.5


async def enrich_context(
    context_items: List[NormalizedContextItem],
    user_id: str,
    query: str,
    user_profile: Dict
) -> List[Dict]:
    """
    Enrich context with pedagogical scores and metadata
    """
    enriched = []
    current_time = datetime.now(timezone.utc)
    
    for item in context_items:
        # Calculate relevance score
        relevance_score = calculate_relevance(item.content, query)
        
        # Calculate recency score
        created_at_timestamp = item.metadata.get("created_at")
        recency_score = calculate_recency_score(created_at_timestamp)
        
        # Calculate engagement score
        engagement_score = calculate_engagement_score(item.source_type, item.metadata)
        
        # Calculate user alignment score
        subject_domain = item.metadata.get("subject_domain", "general")
        user_alignment = calculate_user_alignment(
            subject_domain,
            user_profile,
            item.content
        )
        
        enriched.append({
            **item.__dict__,
            "relevance_score": relevance_score,
            "recency_score": recency_score,
            "engagement_score": engagement_score,
            "pedagogical_level": deduce_pedagogical_level(item.content),
            "subject_domain": item.metadata.get("subject_domain", "general"),
            "user_preference_alignment": user_alignment
        })
    
    return enriched


# ============================================================================
# STEP 4: PEDAGOGICAL FILTERING
# ============================================================================

def filter_context_pedagogical(
    enriched_items: List[Dict],
    user_profile: Dict,
    config: Dict
) -> List[Dict]:
    """
    Filter context according to pedagogical criteria
    """
    relevance_min = config.get("relevance_threshold", 0.3)
    recency_min = config.get("recency_threshold", 0.1)
    
    filtered = []
    level_map = {"beginner": 0, "intermediate": 1, "advanced": 2}
    
    user_level = level_map.get(user_profile.get("pedagogical_level", "intermediate"), 1)
    
    for item in enriched_items:
        # Filter by relevance
        if item.get("relevance_score", 0) < relevance_min:
            continue
        
        # Filter by recency
        if item.get("recency_score", 0) < recency_min:
            continue
        
        # Filter by pedagogical level compatibility
        item_level = level_map.get(item.get("pedagogical_level", "intermediate"), 1)
        level_diff = abs(item_level - user_level)
        
        if level_diff > 1:  # Maximum gap of 1 level
            continue
        
        filtered.append(item)
    
    # Remove duplicates
    filtered = remove_duplicates(filtered)
    
    return filtered


def calculate_cosine_similarity(text1: str, text2: str) -> float:
    """
    Simple cosine similarity based on word overlap
    
    Returns: score between 0 and 1
    """
    if not text1 or not text2:
        return 0.0
    
    words1 = set(text1.lower().split())
    words2 = set(text2.lower().split())
    
    if not words1 or not words2:
        return 0.0
    
    intersection = len(words1 & words2)
    union = len(words1 | words2)
    
    similarity = intersection / union if union > 0 else 0.0
    
    return similarity


def compute_composite_score(item: Dict) -> float:
    """
    Compute composite score for ranking
    
    Weights: relevance 40%, engagement 30%, recency 20%, alignment 10%
    """
    return (
        0.4 * item.get("relevance_score", 0) +
        0.3 * item.get("engagement_score", 0) +
        0.2 * item.get("recency_score", 0) +
        0.1 * item.get("user_preference_alignment", 0)
    )


def remove_duplicates(items: List[Dict], similarity_threshold: float = 0.95) -> List[Dict]:
    """
    Remove duplicate or very similar items
    
    Keeps the item with highest composite score
    """
    unique = []
    
    for item1 in items:
        is_duplicate = False
        
        for i, item2 in enumerate(unique):
            similarity = calculate_cosine_similarity(
                item1.get("content", ""),
                item2.get("content", "")
            )
            
            if similarity > similarity_threshold:
                is_duplicate = True
                
                # Compare composite scores
                score1 = compute_composite_score(item1)
                score2 = compute_composite_score(item2)
                
                if score1 > score2:
                    unique[i] = item1
                
                break
        
        if not is_duplicate:
            unique.append(item1)
    
    return unique


# ============================================================================
# STEP 5: RANKING AND SORTING
# ============================================================================

def rank_context(
    filtered_items: List[Dict],
    weights: Optional[Dict[str, float]] = None,
    diversity_strategy: bool = True,
    max_results: int = 20
) -> List[Dict]:
    """
    Rank context items by composite score
    """
    if weights is None:
        weights = {
            "relevance": 0.4,
            "engagement": 0.3,
            "recency": 0.2,
            "user_alignment": 0.1
        }
    
    # Calculate composite scores
    for item in filtered_items:
        composite_score = (
            weights.get("relevance", 0.4) * item.get("relevance_score", 0) +
            weights.get("engagement", 0.3) * item.get("engagement_score", 0) +
            weights.get("recency", 0.2) * item.get("recency_score", 0) +
            weights.get("user_alignment", 0.1) * item.get("user_preference_alignment", 0)
        )
        item["composite_score"] = composite_score
    
    # Sort by composite score
    ranked = sorted(filtered_items, key=lambda x: x.get("composite_score", 0), reverse=True)
    
    # Apply diversity if requested
    if diversity_strategy:
        ranked = apply_diversity_strategy(ranked)
    
    # Limit results
    ranked = ranked[:max_results]
    
    # Normalize scores and add ranking
    max_score = max([item.get("composite_score", 0) for item in ranked]) if ranked else 1.0
    
    for rank, item in enumerate(ranked, 1):
        item["ranking_position"] = rank
        item["normalized_score"] = (
            item.get("composite_score", 0) / max_score if max_score > 0 else 0
        )
    
    return ranked


def apply_diversity_strategy(items: List[Dict]) -> List[Dict]:
    """
    Apply diversity strategy to avoid source concentration
    """
    # Count source types
    source_counts = {}
    for item in items:
        source_type = item.get("source_type", "unknown")
        source_counts[source_type] = source_counts.get(source_type, 0) + 1
    
    if not source_counts:
        return items
    
    total = len(items)
    max_per_type = max(3, total // len(source_counts))
    
    # Rebalance
    diversified = []
    type_counts = {}
    
    for item in items:
        source_type = item.get("source_type", "unknown")
        type_counts[source_type] = type_counts.get(source_type, 0) + 1
        
        if type_counts[source_type] <= max_per_type:
            diversified.append(item)
    
    return diversified


def format_ranked_output(ranked_items: List[Dict]) -> List[Dict]:
    """
    Format final output for API response
    """
    output = []
    
    for item in ranked_items:
        content = item.get("content", "")
        preview = (content[:300] + "...") if len(content) > 300 else content
        
        output.append({
            "rank": item.get("ranking_position"),
            "source": item.get("source_type"),
            "source_id": item.get("id"),
            "content_preview": preview,
            "full_content": content,
            "metadata": {
                "type": item.get("metadata", {}).get("type", ""),
                "created_at": item.get("metadata", {}).get("created_at"),
                "last_updated": item.get("metadata", {}).get("last_updated"),
                "source_id": item.get("id"),
                "title": item.get("metadata", {}).get("title"),
                "subject_domain": item.get("subject_domain", "general"),
                "source": item.get("metadata", {}).get("source", ""),
                "interaction_count": item.get("metadata", {}).get("interaction_count", 0),
                "answered_questions": item.get("metadata", {}).get("answered_questions", 0)
            },
            "scores": {
                "relevance": round(item.get("relevance_score", 0), 3),
                "engagement": round(item.get("engagement_score", 0), 3),
                "recency": round(item.get("recency_score", 0), 3),
                "user_alignment": round(item.get("user_preference_alignment", 0), 3),
                "composite": round(item.get("composite_score", 0), 3),
                "normalized": round(item.get("normalized_score", 0), 3)
            }
        })
    
    return output


# ============================================================================
# API ROUTES
# ============================================================================

@router.post("/context/retrieve", response_model=List[ContextRetrievalResponse])
async def retrieve_context(
    request: ContextRetrievalRequest,
    user=Depends(get_verified_user),
):
    """
    Retrieve relevant context from multiple sources:
    - Pedagogical documents (RAG)
    - Internal memory (episodic, semantic, procedural, behavioral)
    - Generated summaries
    
    Results are filtered and ranked by pedagogical relevance.
    
    Query Parameters:
    - query: The search query
    - max_results: Maximum number of results (default: 20)
    - memory_types: Filter by memory types (optional)
    - pedagogical_level: User pedagogical level (default: intermediate)
    """
    db = get_db_session()
    try:
        # Build user profile
        user_profile = {
            "pedagogical_level": request.pedagogical_level or "intermediate",
            "interests": [],
            "learning_objectives": request.learning_objectives or []
        }
        
        # STEP 1: Retrieve from multiple sources
        documents = await retrieve_pedagogical_documents(
            user.id, request.query, top_k=5
        )
        
        memories = await retrieve_internal_memory(
            user.id,
            request.query,
            memory_types=request.memory_types,
            limit=10,
            db=db,
            topic=request.topic,
        )

        summaries = await retrieve_generated_summaries(
            user.id, request.query, limit=5, topic=request.topic,
        )
        
        # Generate lightweight summaries when cache is empty
        if not summaries:
            summaries = []
            for item in (documents[:2] + memories[:2]):
                content = item.get("content", "")
                if content:
                    summaries.append({
                        "id": f"generated-{item.get('id', 'unknown')}",
                        "text": await summarize_interactions(content, query=request.query, max_tokens=300),
                        "source_conversation": item.get("id", ""),
                        "created_at": datetime.now(timezone.utc).timestamp(),
                        "summary_score": calculate_relevance(content, request.query)
                    })
        if request.include_source_types:
            source_types = request.include_source_types
            if "pedagogical" not in source_types:
                documents = []
            if "memory" not in source_types:
                memories = []
            if "summary" not in source_types:
                summaries = []
        
        # STEP 2: Normalize
        normalized = normalize_context(documents, memories, summaries)
        
        if not normalized:
            return []
        
        # STEP 3: Enrich with pedagogical metadata
        enriched = await enrich_context(
            normalized, user.id, request.query, user_profile
        )
        
        # STEP 4: Filter pedagogically
        filtered_config = {
            "relevance_threshold": 0.3,
            "recency_threshold": 0.1
        }
        filtered = filter_context_pedagogical(enriched, user_profile, filtered_config)
        
        if not filtered:
            return []
        
        # STEP 5: Rank and sort
        ranked = rank_context(
            filtered,
            weights={
                "relevance": 0.4,
                "engagement": 0.3,
                "recency": 0.2,
                "user_alignment": 0.1
            },
            diversity_strategy=True,
            max_results=request.max_results
        )
        
        # STEP 6: Apply summarization layer
        summarization_config = CONTEXT_RETRIEVAL_CONFIG.get("summarization", {}).copy()
        summarization_config["sliding_window_size"] = request.max_results
        summarized = await apply_summarization_layer(ranked, request.query, {"summarization": summarization_config})
        
        # Format output
        result = format_ranked_output(summarized)
        
        return result
        
    except Exception as exc:
        print(f"Error in context retrieval: {str(exc)}")
        raise HTTPException(status_code=500, detail=f"Context retrieval failed: {str(exc)}")
    finally:
        db.close()


@router.get("/context/stats")
async def get_context_stats(
    user=Depends(get_verified_user),
):
    """
    Get statistics about available context sources for the user
    """
    db = get_db_session()
    try:
        # Count memories by type
        memory_query = db.query(Memory).filter(Memory.user_id == user.id)
        total_memories = memory_query.count()
        
        memory_types_count = {}
        for memory_type in ["episodic", "semantic", "procedural", "behavioral"]:
            count = memory_query.filter(Memory.memory_type == memory_type).count()
            memory_types_count[memory_type] = count

        behavioral_count = memory_types_count.get("behavioral", 0)
        pedagogical_document_count = 0
        try:
            collection = get_or_create_collection(DOCUMENTS_COLLECTION)
            pedagogical_document_count = collection.count()
        except Exception:
            pedagogical_document_count = 0

        return {
            "user_id": user.id,
            "total_memories": total_memories,
            "memory_types": memory_types_count,
            "behavioral_memories": behavioral_count,
            "pedagogical_document_count": pedagogical_document_count,
            "available_sources": ["memory", "pedagogical"],
            "note": "Context stats now include internal memories, behavioral interaction memories, and indexed pedagogical documents."
        }
        
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc))
    finally:
        db.close()

@router.post("/context/index-local-documents")
async def index_local_documents_endpoint(
    user=Depends(get_verified_user)
):
    """
    Manually trigger indexing of local documents into ChromaDB
    
    This endpoint allows administrators to re-index local documents
    if the collection becomes out of sync.
    """
    try:
        indexed_count = index_local_documents_to_chromadb()
        
        return {
            "status": "success",
            "indexed_count": indexed_count,
            "message": f"Indexed {indexed_count} local documents into ChromaDB"
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to index documents: {str(e)}")


def _get_document_metadata(doc_id: str) -> Optional[Dict[str, Any]]:
    """Retrieve raw document metadata from ChromaDB by document id."""
    try:
        from chromadb.api.types import Include
        collection = get_or_create_collection(DOCUMENTS_COLLECTION)
        result = collection.get(ids=[doc_id], include=cast(Include, ["metadatas"]))
        if not result or not result.get("ids"):
            return None
        metadatas = result.get("metadatas", [])
        if metadatas:
            # Convert Metadata object to dict if needed
            metadata = metadatas[0]
            if hasattr(metadata, '__dict__'):
                return dict(metadata)
            elif isinstance(metadata, dict):
                return metadata
            else:
                return None
        return None
    except Exception:
        return None


def _update_document_metadata(doc_id: str, updates: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """Update document metadata in ChromaDB for a given document id."""
    try:
        collection = get_or_create_collection(DOCUMENTS_COLLECTION)
        existing_metadata = _get_document_metadata(doc_id) or {}
        merged_metadata = {**existing_metadata, **updates}
        collection.update(ids=[doc_id], metadatas=[merged_metadata])
        return merged_metadata
    except Exception:
        return None


def _find_behavioral_memory_for_doc(user_id: str, doc_id: str):
    """Find an existing behavioral memory for the same user and document."""
    db = get_db_session()
    try:
        memories = db.query(Memory).filter(
            Memory.user_id == user_id,
            Memory.memory_type == "behavioral"
        ).all()
        for memory in memories:
            metadata = memory.memory_metadata or {}
            if metadata.get("document_id") == doc_id:
                return memory
        return None
    except Exception:
        return None
    finally:
        db.close()


@router.post("/context/documents/{doc_id}/interaction")
async def record_document_interaction(
    doc_id: str,
    request: DocumentInteractionRequest,
    user=Depends(get_verified_user)
):
    """Record a user interaction with a pedagogical document and update engagement metadata."""
    existing_metadata = _get_document_metadata(doc_id)
    if existing_metadata is None:
        raise HTTPException(status_code=404, detail=f"Document '{doc_id}' not found in vector DB")

    interaction_count = int(existing_metadata.get("interaction_count", 0))
    answered_questions = int(existing_metadata.get("answered_questions", 0))

    interaction_count += request.interaction_increment
    answered_questions += request.answered_questions_increment or 0

    update_fields = {
        "interaction_count": interaction_count,
        "answered_questions": answered_questions,
        "last_updated": request.last_updated or datetime.now(timezone.utc).timestamp(),
    }
    if request.interaction_type:
        update_fields["last_interaction_type"] = request.interaction_type

    updated_metadata = _update_document_metadata(doc_id, update_fields)
    if updated_metadata is None:
        raise HTTPException(status_code=500, detail="Failed to update document interaction metadata")

    # Persist or update a behavioral memory for this document interaction.
    db = get_db_session()
    try:
        memory = _find_behavioral_memory_for_doc(user.id, doc_id)
        memory_metadata = {
            "document_id": doc_id,
            "interaction_type": request.interaction_type,
            "interaction_count": interaction_count,
            "answered_questions": answered_questions,
            "last_updated": updated_metadata.get("last_updated")
        }

        if memory is not None:
            setattr(memory, 'content', 
                f"Interaction with document {doc_id}: {request.interaction_type or 'interaction'}"
            )
            setattr(memory, 'memory_metadata', memory_metadata)
            db.add(memory)
        else:
            memory = Memory(
                id=uuid4().hex,
                user_id=user.id,
                memory_type="behavioral",
                content=(
                    f"Interaction with document {doc_id}: {request.interaction_type or 'interaction'}"
                ),
                memory_metadata=memory_metadata,
            )
            db.add(memory)

        db.commit()
        db.refresh(memory)
    except Exception as exc:
        db.rollback()
        print(f"Failed to persist interaction memory for doc {doc_id}: {exc}")
    finally:
        db.close()

    return {
        "status": "success",
        "doc_id": doc_id,
        "interaction_count": interaction_count,
        "answered_questions": answered_questions,
        "last_updated": updated_metadata.get("last_updated")
    }


@router.delete("/context/reset-vector-db")
async def reset_vector_database(
    user=Depends(get_verified_user)
):
    try:
        # Delete the collection
        try:
            _get_chroma_client().delete_collection(name=DOCUMENTS_COLLECTION)
        except Exception:
            pass  # Collection might not exist

        # Recreate the collection
        collection = _get_chroma_client().create_collection(name=DOCUMENTS_COLLECTION)

        return {
            "status": "success",
            "message": f"Reset ChromaDB collection '{DOCUMENTS_COLLECTION}'"
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to reset vector database: {str(e)}")
    
    
@router.get("/context/vector-db-stats")
async def get_vector_db_stats(
    user=Depends(get_verified_user)
):
    """
    Get statistics about the ChromaDB vector database
    """
    try:
        collection = get_or_create_collection(DOCUMENTS_COLLECTION)
        count = collection.count()
        
        return {
            "collection_name": DOCUMENTS_COLLECTION,
            "document_count": count,
            "status": "active"
        }
        
    except Exception as e:
        return {
            "collection_name": DOCUMENTS_COLLECTION,
            "document_count": 0,
            "status": "error",
            "error": str(e)
        }
